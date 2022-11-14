import os
import pandas as pd
from pptx import Presentation, shapes, presentation
import win32com.client
import smtplib
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from email.mime.application import MIMEApplication
from dotenv import load_dotenv

load_dotenv()


## Environment Variables:
APP_GMAIL_ADDRESS = os.getenv("APP_GMAIL_ADDRESS")
APP_GMAIL_PASSWORD = os.getenv("APP_GMAIL_PASSWORD")


## Global Variables:
ROOT_DIRECTORY = r"C:\Users\sputnam\Documents\AWRA\11_7_22_Seminar"
PDH_DIRECTORY = os.path.join(ROOT_DIRECTORY, "PDH")
PDH_FILENAME = "AWRA_PDH_11_7_22_Seminar_{}.pptx"

PDH_EMAIL_CC = ["kearam55@gmail.com", "powersdb@cdmsmith.com", "msreetharan@dewberry.com"]
PDH_EMAIL_SUBJECT = "AWRA-NCRS Certificate of Attendance"
PDH_EMAIL_BODY = """Dear {},\n
Thank you for attending the AWRA-NCRS event titled “Coastal Development in the Nation’s Capital and Beyond.” Please see attached for your certificate of attendance.\n
As a reminder, we will be hosting another evening seminar in January 2023, so please keep an eye out for additional details.\n
In the next few weeks, we will be issuing a call for abstracts for our Annual Water Symposium. The symposium’s theme is Equitable and Resilient Water Management and will take place in April 2023.\n
To stay up-to-date with all of the latest news please visit our website: http://www.awrancrs.org/.\n
Sincerely,
Shane\n
Shane M. Putnam, Ph.D.
AWRA-NCRS Treasurer
"""


## Functions:
def update_save_powerpoint(
    textbox: shapes.autoshape.Shape,
    updated_text: str,
    pdh_template: presentation.Presentation,
    attendee_pdh_pptx_path: str,
) -> None:
    """ """
    paragraph = textbox.text_frame.paragraphs[0]
    paragraph.runs[0].text = updated_text
    pdh_template.save(attendee_pdh_pptx_path)
    return


def replace_powerpoint_w_pdf(powerpoint_path: str, pdf_path: str) -> None:
    """ """
    powerpoint = win32com.client.Dispatch("Powerpoint.Application")
    deck = powerpoint.Presentations.Open(powerpoint_path)
    deck.SaveAs(pdf_path, 32)
    deck.Close()
    os.remove(powerpoint_path)
    return


def create_email(
    sender_email_address: str, cc_email_address: str, to_email_address: str, email_subject: str, email_body: str
) -> MIMEMultipart:
    """ """
    message = MIMEMultipart()
    message["From"] = sender_email_address
    message["Cc"] = cc_email_address
    message["To"] = to_email_address
    message["Subject"] = email_subject
    message.attach(MIMEText(email_body))
    return message


def add_attachment(message: MIMEMultipart, attachment_path: str) -> MIMEMultipart:
    """ """
    attachment_filename = os.path.basename(attachment_path)
    with open(attachment_path, "rb") as f:
        part = MIMEApplication(f.read(), Name=attachment_filename)
    part["Content-Disposition"] = f'attachment; filename="{attachment_filename}"'
    message.attach(part)
    return message


def send_email(
    message: MIMEMultipart, sender_email_address: str, sender_password: str, recipients_email_addresses: list
) -> None:
    """ """
    session = smtplib.SMTP("smtp.gmail.com", 587)
    session.starttls()
    session.login(sender_email_address, sender_password)  # login with mail_id and password
    session.sendmail(sender_email_address, recipients_email_addresses, message.as_string())
    session.quit()
    recipients_str = ", ".join(recipients_email_addresses)
    print(f"Sent to {recipients_str}")
    return None


## Read in the attendance list:
attendees_filename = "november_attendees.xlsx"
attendees_path = os.path.join(ROOT_DIRECTORY, attendees_filename)

attendees = pd.read_excel(attendees_path)

print(f"There were {attendees['attended'].sum()} attendees and they paid ${attendees['amount_paid'].sum()}\n")

print(
    f"\nThe following details are missing:\n\t- First Name: {attendees['first'].isnull().sum()}\n\t- Last Name: "
    f"{attendees['last'].isnull().sum()}\n\t- Email: {attendees['email'].isnull().sum()}\n"
)

attendees = attendees[(attendees["attended"])].copy()

for col in ["first", "last", "email"]:
    attendees[col] = attendees[col].str.strip()


## Read in the certificate template:
pdh_template_filename = "awra_pdh_2022_2023_template.pptx"
pdh_template_path = os.path.join(ROOT_DIRECTORY, pdh_template_filename)

pdh_template = Presentation(pdh_template_path)
pdh_template_slide = pdh_template.slides[0]


## For each attendee make a certificate:
if not os.path.exists(PDH_DIRECTORY):
    os.makedirs(PDH_DIRECTORY)

textbox = [shape for shape in pdh_template_slide.shapes if shape.shape_type == 17 and "First Last" in shape.text][0]

for i, attendee in enumerate(attendees.itertuples()):
    first_name = attendee.first
    last_name = attendee.last
    participant_email = attendee.email

    # Make certificate and convert to PDF:
    attendee_pdh_pptx_path = os.path.join(PDH_DIRECTORY, PDH_FILENAME.format(last_name))
    updated_text = f"{first_name} {last_name}"
    update_save_powerpoint(textbox, updated_text, pdh_template, attendee_pdh_pptx_path)

    attendee_pdh_pdf_path = attendee_pdh_pptx_path.replace(".pptx", ".pdf")
    replace_powerpoint_w_pdf(attendee_pdh_pptx_path, attendee_pdh_pdf_path)

    if i + 1 == attendees.shape[0]:
        powerpoint = win32com.client.Dispatch("Powerpoint.Application")
        powerpoint.Quit()

    # Email certificate
    cc_list = [email for email in PDH_EMAIL_CC if email != participant_email]
    cc_email_address = ",".join(cc_list)
    recipients_email_addresses = cc_list + [participant_email]
    email_body = PDH_EMAIL_BODY.format(first_name)

    message = create_email(APP_GMAIL_ADDRESS, cc_email_address, participant_email, PDH_EMAIL_SUBJECT, email_body)
    message = add_attachment(message, attendee_pdh_pdf_path)
    send_email(message, APP_GMAIL_ADDRESS, APP_GMAIL_PASSWORD, recipients_email_addresses)
